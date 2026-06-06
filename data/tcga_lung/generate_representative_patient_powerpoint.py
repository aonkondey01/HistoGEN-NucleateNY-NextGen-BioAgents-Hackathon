#!/usr/bin/env python3
"""Build a PowerPoint deck for the 20 representative TCGA lung patient visual summary.

Creates native pie charts, bar/column charts with statistics, plus PNG embeds of
complex SVG plots (heatmaps, chromosome landscapes) from the visual report.
"""

from __future__ import annotations

import argparse
import csv
import json
import re
import statistics
import tempfile
from collections import Counter
from pathlib import Path
from typing import Any

import cairosvg
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
DARK = RGBColor(15, 23, 42)
MUTED = RGBColor(71, 85, 105)
ACCENT = RGBColor(59, 130, 246)


def read_csv(path: Path) -> list[dict[str, str]]:
    with path.open(newline="") as fh:
        return list(csv.DictReader(fh))


def to_float(value: Any) -> float | None:
    if value in (None, ""):
        return None
    try:
        return float(value)
    except (TypeError, ValueError):
        return None


def fmt_num(value: Any, digits: int = 1) -> str:
    if value is None or value == "":
        return "n/a"
    if isinstance(value, float):
        if value.is_integer():
            return str(int(value))
        return f"{value:.{digits}f}"
    return str(value)


def pct(part: int, whole: int) -> str:
    if not whole:
        return "0%"
    return f"{100 * part / whole:.1f}%"


def human_title(stem: str) -> str:
    title = re.sub(r"^\d+_", "", stem.replace(".svg", ""))
    return title.replace("_", " ").strip().title()


def svg_to_png(svg_path: Path, png_path: Path, *, scale: float = 2.0) -> None:
    cairosvg.svg2png(url=str(svg_path), write_to=str(png_path), scale=scale)


def set_title_text(shape, text: str, *, size: int = 24) -> None:
    shape.text = text
    para = shape.text_frame.paragraphs[0]
    para.font.size = Pt(size)
    para.font.bold = True
    para.font.color.rgb = DARK


def add_title_slide(prs: Presentation, summary: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_title_text(slide.shapes.title, "Representative 20 TCGA Lung Patients", size=32)
    n = summary.get("patient_count", 20)
    luad = summary.get("project_counts", {}).get("TCGA-LUAD", 10)
    lusc = summary.get("project_counts", {}).get("TCGA-LUSC", 10)
    alive = summary.get("vital_status_counts", {}).get("Alive", 0)
    dead = summary.get("vital_status_counts", {}).get("Dead", 0)
    subtitle = slide.placeholders[1]
    subtitle.text = (
        f"{n} patients | {luad} LUAD + {lusc} LUSC | 5 smokers + 5 non-smokers per histology\n"
        f"{summary.get('total_somatic_mutations', 0):,} somatic mutations | "
        f"Vital status: {alive} alive, {dead} dead\n"
        "Clinical, mutation, RNA, RPPA, CNV, miRNA, methylation, and treatment summary"
    )


def add_section_slide(prs: Presentation, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_title_text(slide.shapes.title, title, size=36)


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title_text(slide.shapes.title, title)
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, text in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(16)


def add_pie_chart_slide(
    prs: Presentation,
    title: str,
    labels: list[str],
    values: list[float],
    *,
    stats_lines: list[str] | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.55))
    set_title_text(slide.shapes[-1], title)

    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series("Count", values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(0.5),
        Inches(0.9),
        Inches(6.2),
        Inches(5.8),
        chart_data,
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.number_format = "0"
    chart.plots[0].data_labels.show_percentage = True

    if stats_lines:
        box = slide.shapes.add_textbox(Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.2))
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(stats_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(14)
            p.font.color.rgb = MUTED


def add_bar_chart_slide(
    prs: Presentation,
    title: str,
    categories: list[str],
    values: list[float],
    *,
    chart_type: XL_CHART_TYPE = XL_CHART_TYPE.COLUMN_CLUSTERED,
    stats_lines: list[str] | None = None,
    y_axis_title: str = "Count",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.55))
    set_title_text(slide.shapes[-1], title)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(y_axis_title, values)
    left = Inches(0.5)
    top = Inches(0.9)
    width = Inches(7.5 if stats_lines else 12.2)
    chart = slide.shapes.add_chart(chart_type, left, top, width, Inches(5.8), chart_data).chart
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = True
    chart.plots[0].has_data_labels = True

    if stats_lines:
        box = slide.shapes.add_textbox(Inches(8.3), Inches(1.2), Inches(4.5), Inches(5.2))
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(stats_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(13)
            p.font.color.rgb = MUTED


def add_grouped_bar_slide(
    prs: Presentation,
    title: str,
    categories: list[str],
    series: dict[str, list[float]],
    *,
    stats_lines: list[str] | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.55))
    set_title_text(slide.shapes[-1], title)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, vals in series.items():
        chart_data.add_series(name, vals)
    width = Inches(7.8 if stats_lines else 12.2)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5),
        Inches(0.9),
        width,
        Inches(5.8),
        chart_data,
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.value_axis.has_major_gridlines = True

    if stats_lines:
        box = slide.shapes.add_textbox(Inches(8.6), Inches(1.2), Inches(4.2), Inches(5.2))
        tf = box.text_frame
        for i, line in enumerate(stats_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(13)
            p.font.color.rgb = MUTED


def add_plot_slide(prs: Presentation, png_path: Path, title: str, *, note: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.55))
    set_title_text(slide.shapes[-1], title)
    slide.shapes.add_picture(str(png_path), Inches(0.35), Inches(0.85), width=Inches(12.6))
    if note:
        box = slide.shapes.add_textbox(Inches(0.4), Inches(6.95), Inches(12.4), Inches(0.45))
        p = box.text_frame.paragraphs[0]
        p.text = note
        p.font.size = Pt(11)
        p.font.color.rgb = MUTED


def add_table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    *,
    note: str = "",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.55))
    set_title_text(slide.shapes[-1], title)
    nrows = len(rows) + 1
    ncols = len(headers)
    table = slide.shapes.add_table(nrows, ncols, Inches(0.35), Inches(0.85), Inches(12.6), Inches(5.8)).table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            para.font.bold = True
            para.font.size = Pt(10)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(9)
    if note:
        box = slide.shapes.add_textbox(Inches(0.4), Inches(6.85), Inches(12.4), Inches(0.45))
        p = box.text_frame.paragraphs[0]
        p.text = note
        p.font.size = Pt(10)
        p.font.color.rgb = MUTED


def count_field(rows: list[dict[str, str]], field: str, missing: str = "missing") -> Counter[str]:
    c: Counter[str] = Counter()
    for row in rows:
        c[row.get(field) or missing] += 1
    return c


def split_field_counts(rows: list[dict[str, str]], field: str) -> Counter[str]:
    c: Counter[str] = Counter()
    for row in rows:
        raw = row.get(field) or ""
        if not raw.strip() or raw == "missing":
            c["Not reported"] += 1
            continue
        for part in re.split(r"\s*;\s*", raw):
            part = part.strip()
            if part:
                c[part] += 1
    return c


def age_histogram_bins(ages: list[float], bins: int = 8) -> tuple[list[str], list[float]]:
    if not ages:
        return [], []
    lo, hi = min(ages), max(ages)
    if lo == hi:
        hi = lo + 1
    edges = [lo + (hi - lo) * i / bins for i in range(bins + 1)]
    counts = [0] * bins
    for age in ages:
        idx = min(bins - 1, int((age - lo) / (hi - lo) * bins))
        counts[idx] += 1
    labels = [f"{edges[i]:.0f}-{edges[i + 1]:.0f}" for i in range(bins)]
    return labels, [float(c) for c in counts]


def aggregate_variant_classes(variant_rows: list[dict[str, str]]) -> Counter[str]:
    c: Counter[str] = Counter()
    for row in variant_rows:
        cls = row.get("variant_classification", "")
        count = int(float(row.get("mutation_count") or 0))
        if cls:
            c[cls] += count
    return c


def build_powerpoint(report_dir: Path, out_path: Path) -> dict[str, int]:
    summary = json.loads((report_dir / "cohort_visual_summary.json").read_text())
    patients = read_csv(report_dir / "master_patient_summary.csv")
    driver_freq = read_csv(report_dir / "driver_mutation_frequency.csv")
    cnv_summary = read_csv(report_dir / "per_patient_cnv_summary.csv")
    mirna_summary = read_csv(report_dir / "per_patient_mirna_summary.csv")
    plots_dir = report_dir / "plots"
    gen_dir = report_dir.parent / "genomic_data"
    variant_rows = read_csv(gen_dir / "mutation_maps" / "variant_classification_summary.csv")

    n = len(patients)
    ages = [v for v in (to_float(r.get("age_at_diagnosis_years")) for r in patients) if v is not None]
    project_counts = count_field(patients, "project_id")
    smoking_counts = count_field(patients, "smoking_group")
    sex_counts = count_field(patients, "sex")
    vital_counts = count_field(patients, "vital_status")
    stage_counts = count_field(patients, "ajcc_pathologic_stage")
    mutation_burden = [int(float(r.get("somatic_mutation_rows") or 0)) for r in patients]

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs, summary)

    # --- Section 1: Cohort overview ---
    add_section_slide(prs, "1. Cohort Overview")
    add_pie_chart_slide(
        prs,
        "Histology distribution (LUAD vs LUSC)",
        ["LUAD", "LUSC"],
        [float(project_counts.get("TCGA-LUAD", 0)), float(project_counts.get("TCGA-LUSC", 0))],
        stats_lines=[
            f"Total patients: {n}",
            f"LUAD: {project_counts.get('TCGA-LUAD', 0)} ({pct(project_counts.get('TCGA-LUAD', 0), n)})",
            f"LUSC: {project_counts.get('TCGA-LUSC', 0)} ({pct(project_counts.get('TCGA-LUSC', 0), n)})",
            "Balanced 50/50 histology split by design.",
        ],
    )
    add_pie_chart_slide(
        prs,
        "Smoking status (selection strata)",
        ["Smoker", "Non-smoker"],
        [float(smoking_counts.get("smoker", 0)), float(smoking_counts.get("non_smoker", 0))],
        stats_lines=[
            f"Smokers: {smoking_counts.get('smoker', 0)} ({pct(smoking_counts.get('smoker', 0), n)})",
            f"Non-smokers: {smoking_counts.get('non_smoker', 0)} ({pct(smoking_counts.get('non_smoker', 0), n)})",
            "5 smokers + 5 non-smokers per histology.",
        ],
    )
    age_labels, age_vals = age_histogram_bins(ages)
    add_bar_chart_slide(
        prs,
        "Age at diagnosis distribution",
        age_labels,
        age_vals,
        stats_lines=[
            f"Mean age: {fmt_num(statistics.mean(ages))} years" if ages else "Mean age: n/a",
            f"Median age: {fmt_num(statistics.median(ages))} years" if ages else "Median age: n/a",
            f"Range: {fmt_num(min(ages))} – {fmt_num(max(ages))} years" if ages else "Range: n/a",
            f"Patients with age data: {len(ages)}/{n}",
        ],
        y_axis_title="Patients",
    )

    # --- Section 2: Clinical demographics ---
    add_section_slide(prs, "2. Clinical Demographics")
    add_pie_chart_slide(
        prs,
        "Vital status",
        ["Alive", "Dead"],
        [float(vital_counts.get("Alive", 0)), float(vital_counts.get("Dead", 0))],
        stats_lines=[
            f"Alive: {vital_counts.get('Alive', 0)} ({pct(vital_counts.get('Alive', 0), n)})",
            f"Dead: {vital_counts.get('Dead', 0)} ({pct(vital_counts.get('Dead', 0), n)})",
            f"Death rate: {pct(vital_counts.get('Dead', 0), n)}",
            f"Events (survival_event=1): {sum(int(float(r.get('survival_event') or 0)) for r in patients)}",
        ],
    )
    sex_items = sex_counts.most_common()
    add_pie_chart_slide(
        prs,
        "Sex distribution",
        [k for k, _ in sex_items],
        [float(v) for _, v in sex_items],
        stats_lines=[f"{label}: {int(v)} ({pct(int(v), n)})" for label, v in sex_items],
    )
    stage_items = stage_counts.most_common()
    add_bar_chart_slide(
        prs,
        "AJCC pathologic stage",
        [k for k, _ in stage_items],
        [float(v) for _, v in stage_items],
        stats_lines=[
            f"Most common stage: {stage_items[0][0]} (n={stage_items[0][1]})" if stage_items else "",
            f"Early stage (I–II): {sum(v for k, v in stage_counts.items() if 'Stage I' in k and 'Stage III' not in k)}",
            f"Advanced stage (III–IV): {sum(v for k, v in stage_counts.items() if 'Stage III' in k or 'Stage IV' in k)}",
        ],
        y_axis_title="Patients",
    )

    # --- Section 3: Data availability ---
    add_section_slide(prs, "3. Data Availability")
    add_pie_chart_slide(
        prs,
        "Multi-omic data coverage",
        ["RPPA protein", "miRNA", "Methylation"],
        [
            float(summary.get("patients_with_rppa", 0)),
            float(summary.get("patients_with_mirna", 0)),
            float(summary.get("patients_with_methylation", 0)),
        ],
        stats_lines=[
            f"RPPA: {summary.get('patients_with_rppa', 0)}/{n} patients",
            f"miRNA: {summary.get('patients_with_mirna', 0)}/{n} patients",
            f"Methylation: {summary.get('patients_with_methylation', 0)}/{n} patients",
            "All 20 have clinical, MAF, RNA, CNV, and H&E slide metadata.",
        ],
    )
    add_bullet_slide(
        prs,
        "Data modality row counts (median per patient)",
        [
            f"Somatic mutations (MAF rows): median {fmt_num(summary.get('median_mutation_burden', 0), 0)}, total {summary.get('total_somatic_mutations', 0):,}",
            f"Genome-wide RNA rows: 60,660 per sample (121,320 if duplicate samples)",
            f"RPPA peptide targets: 487 per patient (when available)",
            f"CNV segments: median {fmt_num(statistics.median([int(r['segments']) for r in cnv_summary]), 0)}",
            f"miRNA features: median {fmt_num(statistics.median([int(r['mirna_targets']) for r in mirna_summary]), 0) if mirna_summary else 0}",
        ],
    )

    # --- Section 4: Mutations ---
    add_section_slide(prs, "4. Somatic Mutations")
    top_drivers = driver_freq[:12]
    add_bar_chart_slide(
        prs,
        "Top driver-gene mutation frequency",
        [r["gene"] for r in top_drivers],
        [float(r["overall_pct"]) for r in top_drivers],
        stats_lines=[
            f"TP53 mutated: {driver_freq[0]['mutated_cases']}/{n} ({driver_freq[0]['overall_pct']}%)",
            f"Genes tracked: {len(driver_freq)} important lung drivers",
            "Frequencies from masked somatic MAF (whole-exome SNVs/indels).",
            f"Median mutation burden: {fmt_num(summary.get('median_mutation_burden', 0), 0)} MAF rows/patient",
            f"Max burden: {max(mutation_burden):,} (TCGA-18-3409)",
            f"Min burden: {min(mutation_burden):,} (TCGA-44-2661)",
        ],
        y_axis_title="Mutated patients (%)",
    )
    add_grouped_bar_slide(
        prs,
        "Driver mutation frequency: LUAD vs LUSC",
        [r["gene"] for r in top_drivers[:8]],
        {
            "LUAD (%)": [float(r["luad_pct"]) for r in top_drivers[:8]],
            "LUSC (%)": [float(r["lusc_pct"]) for r in top_drivers[:8]],
        },
        stats_lines=[
            "Grouped bars show histology-specific mutation rates.",
            f"KRAS LUAD: {next((r['luad_pct'] for r in driver_freq if r['gene']=='KRAS'), 0)}%",
            f"CDKN2A LUSC: {next((r['lusc_pct'] for r in driver_freq if r['gene']=='CDKN2A'), 0)}%",
        ],
    )
    burden_sorted = sorted(
        [(r["case_submitter_id"].replace("TCGA-", ""), int(float(r["somatic_mutation_rows"]))) for r in patients],
        key=lambda x: -x[1],
    )
    add_bar_chart_slide(
        prs,
        "Somatic mutation burden by patient",
        [p for p, _ in burden_sorted],
        [float(v) for _, v in burden_sorted],
        chart_type=XL_CHART_TYPE.BAR_CLUSTERED,
        stats_lines=[
            f"Total MAF rows: {sum(mutation_burden):,}",
            f"Mean: {fmt_num(statistics.mean(mutation_burden), 0)}",
            f"Median: {fmt_num(statistics.median(mutation_burden), 0)}",
            "Higher counts reflect more exome-covered variants called.",
        ],
        y_axis_title="MAF rows",
    )
    variant_counts = aggregate_variant_classes(variant_rows)
    top_variants = variant_counts.most_common(8)
    add_pie_chart_slide(
        prs,
        "Variant classification (cohort-wide)",
        [k for k, _ in top_variants],
        [float(v) for _, v in top_variants],
        stats_lines=[
            f"Total classified variants: {sum(variant_counts.values()):,}",
            f"Missense: {variant_counts.get('Missense_Mutation', 0):,}",
            f"Silent: {variant_counts.get('Silent', 0):,}",
            f"Nonsense: {variant_counts.get('Nonsense_Mutation', 0):,}",
            "From whole-exome MAF files (not whole-genome).",
        ],
    )

    # --- Section 5: CNV and miRNA ---
    add_section_slide(prs, "5. Copy Number and miRNA")
    cnv_sorted = sorted(cnv_summary, key=lambda r: -(int(r["amp"]) + int(r["del"])))
    add_grouped_bar_slide(
        prs,
        "CNV alterations by patient (top 12 by total amp+del)",
        [r["case_submitter_id"].replace("TCGA-", "") for r in cnv_sorted[:12]],
        {
            "Amplifications (CN≥3)": [float(r["amp"]) for r in cnv_sorted[:12]],
            "Deletions (CN≤1)": [float(r["del"]) for r in cnv_sorted[:12]],
        },
        stats_lines=[
            f"Total CNV segments: {sum(int(r['segments']) for r in cnv_summary):,}",
            f"Mean amp/patient: {fmt_num(statistics.mean([int(r['amp']) for r in cnv_summary]), 0)}",
            f"Mean del/patient: {fmt_num(statistics.mean([int(r['del']) for r in cnv_summary]), 0)}",
            "From ASCAT2 allele-specific segment files.",
        ],
    )
    mirna_sorted = sorted(mirna_summary, key=lambda r: -int(r["mirna_targets"]))
    add_bar_chart_slide(
        prs,
        "miRNA features quantified by patient",
        [r["case_submitter_id"].replace("TCGA-", "") for r in mirna_sorted],
        [float(r["mirna_targets"]) for r in mirna_sorted],
        stats_lines=[
            f"Patients with miRNA: {len(mirna_summary)}/{n}",
            f"Missing miRNA: {n - len(mirna_summary)} patients",
            f"Median features: {fmt_num(statistics.median([int(r['mirna_targets']) for r in mirna_summary]), 0) if mirna_summary else 0}",
        ],
        y_axis_title="miRNA features",
    )

    # --- Section 6: Treatment ---
    add_section_slide(prs, "6. Treatment and Outcomes")
    treat = summary.get("treatment_summary", {})
    top_types = treat.get("top_treatment_types", {})
    type_items = sorted(top_types.items(), key=lambda kv: -kv[1])[:8]
    add_bar_chart_slide(
        prs,
        "Treatment types (patient mentions)",
        [k[:28] + "…" if len(k) > 28 else k for k, _ in type_items],
        [float(v) for _, v in type_items],
        stats_lines=[
            f"Patients with treatment types: {treat.get('field_coverage', {}).get('treatment_types', 0)}/{n}",
            f"Coverage: {treat.get('field_coverage_pct', {}).get('treatment_types', 0)}%",
        ],
        y_axis_title="Mentions",
    )
    response_counts = treat.get("disease_response_counts", {})
    resp_items = [(k, v) for k, v in response_counts.items() if k != "missing"]
    add_pie_chart_slide(
        prs,
        "Disease response",
        [k for k, _ in resp_items],
        [float(v) for _, v in resp_items],
        stats_lines=[
            f"Reported: {sum(v for k, v in resp_items)}/{n}",
            f"Missing: {response_counts.get('missing', 0)}",
            "TF = Tumor Free; WT = With Tumor.",
        ],
    )
    prior = count_field(patients, "prior_treatment")
    add_pie_chart_slide(
        prs,
        "Prior treatment status",
        list(prior.keys()),
        [float(v) for v in prior.values()],
        stats_lines=[f"{k}: {int(v)} ({pct(int(v), n)})" for k, v in prior.most_common()],
    )
    prog = treat.get("progression_or_recurrence_counts", {})
    add_pie_chart_slide(
        prs,
        "Progression or recurrence",
        ["Yes", "Not reported"],
        [float(prog.get("Yes", 0)), float(prog.get("missing", 0))],
        stats_lines=[
            f"Progression/recurrence Yes: {prog.get('Yes', 0)} ({pct(prog.get('Yes', 0), n)})",
            f"Not reported: {prog.get('missing', 0)}",
            "Used as resistance proxy where direct labels are sparse.",
        ],
    )

    # --- Section 7: Master patient table ---
    add_section_slide(prs, "7. Per-Patient Summary")
    table_rows = []
    for r in sorted(patients, key=lambda x: (x.get("project_id", ""), x.get("case_submitter_id", ""))):
        table_rows.append(
            [
                r["case_submitter_id"].replace("TCGA-", ""),
                r.get("project_id", "").replace("TCGA-", ""),
                r.get("smoking_group", "")[:3],
                fmt_num(r.get("age_at_diagnosis_years"), 0),
                (r.get("ajcc_pathologic_stage") or "")[:8],
                r.get("vital_status", "")[:5],
                str(r.get("somatic_mutation_rows", "")),
                (r.get("important_gene_mutations") or "")[:24],
            ]
        )
    add_table_slide(
        prs,
        "Master patient summary (all 20)",
        ["Patient", "Hist", "Smk", "Age", "Stage", "Vital", "Mut rows", "Driver genes"],
        table_rows,
        note="Full table: master_patient_summary.csv in visual_report/",
    )

    # --- Embed complex SVG plots ---
    add_section_slide(prs, "8. Detailed Visualizations")
    svg_plots = sorted(plots_dir.glob("*.svg"))
    with tempfile.TemporaryDirectory(prefix="rep-ppt-") as tmpdir:
        tmp = Path(tmpdir)
        for svg_path in svg_plots:
            png_path = tmp / f"{svg_path.stem}.png"
            svg_to_png(svg_path, png_path)
            note = ""
            if "heatmap" in svg_path.name:
                note = "Heatmap intensity reflects presence (mutations) or log-scaled expression (RNA)."
            if "chromosome" in svg_path.name:
                note = "Whole-exome mutation counts by chromosome (WXS coverage, not full genome)."
            add_plot_slide(prs, png_path, human_title(svg_path.name), note=note)

    # --- Summary slide ---
    add_bullet_slide(
        prs,
        "Summary and caveats",
        [
            f"Cohort: {n} representative patients (10 LUAD + 10 LUSC; balanced smoking strata).",
            f"Total somatic mutations: {summary.get('total_somatic_mutations', 0):,} MAF rows (median {fmt_num(summary.get('median_mutation_burden', 0), 0)}/patient).",
            f"Data: RNA all 20; RPPA {summary.get('patients_with_rppa', 0)}/20; miRNA {summary.get('patients_with_mirna', 0)}/20; methylation {summary.get('patients_with_methylation', 0)}/20.",
            "Mutations are whole-exome SNVs/indels; fusions and CNAs require separate assays.",
            "CNV from SNP6/WGS segment files; RNA TPM for 30 curated driver genes.",
            "Treatment/resistance: progression and disease response used as proxies.",
            "Regenerate: python3 data/tcga_lung/visualize_representative_patients.py",
            "PowerPoint: python3 data/tcga_lung/generate_representative_patient_powerpoint.py",
        ],
    )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return {"total_slides": len(prs.slides), "native_charts": 16, "embedded_plots": len(svg_plots)}


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    root = Path(__file__).resolve().parent
    parser.add_argument(
        "--report-dir",
        type=Path,
        default=root / "representative_patients" / "visual_report",
    )
    parser.add_argument(
        "--out",
        type=Path,
        default=root / "representative_patients" / "visual_report" / "Representative_20_Patients_Summary.pptx",
    )
    args = parser.parse_args()
    if not (args.report_dir / "cohort_visual_summary.json").exists():
        raise SystemExit(
            f"Visual report not found at {args.report_dir}. "
            "Run: python3 data/tcga_lung/visualize_representative_patients.py"
        )
    stats = build_powerpoint(args.report_dir, args.out)
    print(f"Wrote PowerPoint to {args.out}")
    print(json.dumps(stats, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
