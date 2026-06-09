#!/usr/bin/env python3
"""Build a PowerPoint deck from the TCGA lung visual report SVG plots."""

from __future__ import annotations

import argparse
import io
import json
import re
import tempfile
from pathlib import Path

import cairosvg
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

SECTIONS: list[tuple[str, list[str]]] = [
    (
        "Cohort overview",
        [
            "01_cohort_size_and_availability",
            "02_luad_vs_lusc_patient_counts",
            "03_age_distribution",
        ],
    ),
    (
        "Demographics and clinical variables",
        [
            "04_sex_distribution",
            "05_race_distribution",
            "06_ethnicity_distribution",
            "07_smoking_history",
            "08_ajcc_pathologic_stage",
            "09_vital_status",
        ],
    ),
    (
        "Driver mutations",
        ["10_driver_mutation_frequency_luad_lusc"],
    ),
    (
        "Clinical differences by mutation status",
        [
            "11_clinical_death_rate_by_mutation_status",
            "12_clinical_median_age_by_mutation_status",
            "13_clinical_advanced_stage_by_mutation_status",
        ],
    ),
    (
        "Exploratory survival associations",
        [
            "14_mutation_survival_pvalues",
            "15_rna_expression_survival_pvalues",
            "16_protein_expression_survival_pvalues",
        ],
    ),
    (
        "Treatment history and resistance proxies",
        [
            "17_treatment_field_coverage",
            "18_prior_treatment_status",
            "19_treatment_types",
            "20_top_therapeutic_agents",
            "21_disease_response",
            "22_progression_or_recurrence",
            "23_treatment_outcomes",
            "24_resistance_proxies",
        ],
    ),
]


def load_significant_km_sections(plots_dir: Path) -> list[tuple[str, list[str]]]:
    manifest_path = plots_dir / "significant_km_manifest.json"
    if not manifest_path.exists():
        return []
    manifest = json.loads(manifest_path.read_text())
    stems = []
    for rel in manifest.get("plot_files", []):
        rel_path = Path(rel)
        if rel_path.parent.name == "significant_km":
            stems.append(f"significant_km/{rel_path.stem}")
        else:
            stems.append(rel_path.stem)
    if not stems:
        return []
    alpha = manifest.get("significance_alpha", 0.05)
    return [(f"Significant Kaplan-Meier curves (p < {alpha})", stems)]


def human_title(stem: str) -> str:
    title = re.sub(r"^(significant_km/)?", "", stem)
    title = re.sub(r"^\d+_", "", title)
    title = title.replace("_", " ")
    return title[0].upper() + title[1:] if title else stem


def resolve_plot_svg(plots_dir: Path, stem: str) -> Path:
    if stem.startswith("significant_km/"):
        return plots_dir / f"{stem}.svg"
    return plots_dir / f"{stem}.svg"


def svg_to_png(svg_path: Path, png_path: Path, *, scale: float = 2.0) -> None:
    cairosvg.svg2png(
        url=str(svg_path),
        write_to=str(png_path),
        scale=scale,
    )


def add_title_slide(prs: Presentation, summary: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "TCGA Lung Cancer Visual Summary"
    subtitle = slide.placeholders[1]
    subtitle.text = (
        f"956 patients | 1,053 diagnostic slides | "
        f"LUAD {summary['project_counts'].get('TCGA-LUAD', 0)} | "
        f"LUSC {summary['project_counts'].get('TCGA-LUSC', 0)}\n"
        "Important-gene mutations, RNA expression, RPPA protein expression, "
        "and exploratory survival associations"
    )


def add_section_slide(prs: Presentation, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = title


def add_plot_slide(
    prs: Presentation,
    png_path: Path,
    title: str,
    *,
    note: str = "",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.6)).text_frame.text = title
    title_box = slide.shapes[-1]
    title_para = title_box.text_frame.paragraphs[0]
    title_para.font.size = Pt(24)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(15, 23, 42)

    slide.shapes.add_picture(str(png_path), Inches(0.35), Inches(0.85), width=Inches(12.6))

    if note:
        note_box = slide.shapes.add_textbox(Inches(0.4), Inches(6.95), Inches(12.4), Inches(0.45))
        note_frame = note_box.text_frame
        note_frame.text = note
        note_para = note_frame.paragraphs[0]
        note_para.font.size = Pt(11)
        note_para.font.color.rgb = RGBColor(71, 85, 105)


def add_summary_slide(prs: Presentation, summary: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
    slide.shapes.title.text = "Summary and caveats"
    body = slide.placeholders[1].text_frame
    body.clear()
    bullets = [
        f"Patients/cases: {summary['patient_count']}",
        f"Diagnostic slides: {summary['slide_count']}",
        f"Mutation-analyzable patients: {summary['mutation_analyzable_patients']}",
        f"Expression-analyzable patients: {summary['expression_analyzable_patients']}",
        f"Important mutation rows: {summary['important_mutation_rows']}",
        f"Important RNA expression rows: {summary['important_rna_expression_rows']}",
        f"Important RPPA protein rows: {summary['important_protein_expression_rows']}",
        "Survival analyses are exploratory and unadjusted.",
        "RNA/protein survival uses median splits.",
        "Treatment history is partially available; direct drug-resistance labels are sparse.",
        "Resistance proxies include progression/recurrence and progressive disease outcomes.",
        "Mutation MAF files capture SNVs/indels; fusions and CNAs need separate assays.",
    ]
    for i, text in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(16)


def build_powerpoint(
    plots_dir: Path,
    summary_path: Path,
    out_path: Path,
    *,
    include_missing: bool = True,
) -> dict[str, int]:
    summary = json.loads(summary_path.read_text())
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, summary)

    plotted = 0
    missing = 0
    sections = SECTIONS + load_significant_km_sections(plots_dir)
    with tempfile.TemporaryDirectory(prefix="tcga-lung-ppt-") as tmpdir:
        tmp = Path(tmpdir)
        for section_title, stems in sections:
            add_section_slide(prs, section_title)
            for stem in stems:
                svg_path = resolve_plot_svg(plots_dir, stem)
                if not svg_path.exists():
                    if include_missing:
                        missing += 1
                    continue
                png_path = tmp / f"{stem.replace('/', '_')}.png"
                svg_to_png(svg_path, png_path)
                add_plot_slide(prs, png_path, human_title(stem))
                plotted += 1

        # Include any extra non-KM plot files not listed explicitly.
        known = {stem for _, stems in sections for stem in stems}
        extras = sorted(
            p.relative_to(plots_dir).with_suffix("").as_posix()
            for p in plots_dir.glob("*.svg")
            if p.relative_to(plots_dir).with_suffix("").as_posix() not in known
            and "survival_km" not in p.name
        )
        if extras:
            add_section_slide(prs, "Additional plots")
            for stem in extras:
                svg_path = resolve_plot_svg(plots_dir, stem)
                png_path = tmp / f"{stem.replace('/', '_')}.png"
                svg_to_png(svg_path, png_path)
                add_plot_slide(prs, png_path, human_title(stem))
                plotted += 1

    add_summary_slide(prs, summary)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return {"slides_with_plots": plotted, "missing_plots": missing, "total_slides": len(prs.slides)}


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    default_root = Path(__file__).resolve().parent / "important_lung_genes" / "visual_report"
    parser.add_argument("--report-dir", type=Path, default=default_root)
    parser.add_argument(
        "--out",
        type=Path,
        default=default_root / "TCGA_Lung_Visual_Summary.pptx",
    )
    parser.add_argument(
        "--significant-km-out",
        type=Path,
        default=default_root / "TCGA_Lung_Significant_KM.pptx",
    )
    args = parser.parse_args()

    plots_dir = args.report_dir / "plots"
    summary_path = args.report_dir / "cohort_visual_summary.json"
    if not plots_dir.exists():
        raise SystemExit(f"Plots directory not found: {plots_dir}")
    if not summary_path.exists():
        raise SystemExit(f"Summary JSON not found: {summary_path}")

    stats = build_powerpoint(plots_dir, summary_path, args.out)
    print(f"Wrote PowerPoint to {args.out}")
    print(json.dumps(stats, indent=2))

    manifest_path = plots_dir / "significant_km_manifest.json"
    if manifest_path.exists():
        manifest = json.loads(manifest_path.read_text())
        km_sections = load_significant_km_sections(plots_dir)
        if km_sections:
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            add_title_slide(prs, json.loads(summary_path.read_text()))
            with tempfile.TemporaryDirectory(prefix="tcga-lung-km-") as tmpdir:
                tmp = Path(tmpdir)
                for section_title, stems in km_sections:
                    add_section_slide(prs, section_title)
                    for stem in stems:
                        svg_path = resolve_plot_svg(plots_dir, stem)
                        png_path = tmp / f"{stem.replace('/', '_')}.png"
                        svg_to_png(svg_path, png_path)
                        add_plot_slide(prs, png_path, human_title(stem))
            args.significant_km_out.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(args.significant_km_out))
            print(f"Wrote significant KM PowerPoint to {args.significant_km_out}")
            print(json.dumps(manifest, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
