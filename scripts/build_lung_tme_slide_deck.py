#!/usr/bin/env python3
"""Build the PEAT-Nucleate lung cancer / TME PowerPoint deck."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs" / "slides"
OUT_PATH = OUT_DIR / "PEAT-Nucleate-Lung-TME-Deck.pptx"

# Brand palette
NAVY = RGBColor(0, 51, 102)
TEAL = RGBColor(0, 128, 128)
SLATE = RGBColor(55, 65, 81)
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(245, 247, 250)
ACCENT_ORANGE = RGBColor(204, 85, 0)


def _set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(9), Inches(0.35))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(9)
    p.font.color.rgb = SLATE
    p.alignment = PP_ALIGN.RIGHT


def _add_title_bar(slide, title: str, subtitle: str | None = None) -> None:
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.15))  # rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.2), Inches(8.9), Inches(0.65))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.78), Inches(8.9), Inches(0.35))
        stf = sub_box.text_frame
        stf.text = subtitle
        sp = stf.paragraphs[0]
        sp.font.size = Pt(14)
        sp.font.color.rgb = RGBColor(200, 220, 235)


def _add_bullets(
    slide,
    items: list[str],
    left: float = 0.65,
    top: float = 1.45,
    width: float = 8.7,
    height: float = 5.4,
    font_size: int = 20,
    level_indent: bool = False,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if level_indent and item.startswith("  "):
            p.level = 1
            p.text = item.strip()
            p.font.size = Pt(font_size - 2)
        else:
            p.text = item.lstrip()
            p.font.size = Pt(font_size)
        p.font.color.rgb = SLATE
        p.space_after = Pt(10)


def _add_notes(slide, notes: str) -> None:
    slide.notes_slide.notes_text_frame.text = notes


def _add_table_slide(slide, headers: list[str], rows: list[list[str]], col_widths: list[float]) -> None:
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.45), Inches(1.5), Inches(9.1), Inches(5.2))
    table = table_shape.table

    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = SLATE


def build_deck() -> Path:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    footer = "PEAT-Nucleate BioHack 2026  |  Lung cancer TME mapping"

    # --- Slide 1: Title ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, NAVY)
    tbox = s.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(2.2))
    tf = tbox.text_frame
    tf.text = "Same Mutation.\nDifferent Outcome."
    for i, line in enumerate(["Same Mutation.", "Different Outcome."]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE
    sub = s.shapes.add_textbox(Inches(0.72), Inches(4.0), Inches(8.5), Inches(1.2))
    stf = sub.text_frame
    stf.text = (
        "Mapping the tumor microenvironment in lung cancer\n"
        "when targeted therapies fail despite actionable mutations"
    )
    for p in stf.paragraphs:
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(190, 215, 230)
    cred = s.shapes.add_textbox(Inches(0.72), Inches(5.5), Inches(8.5), Inches(0.8))
    ctf = cred.text_frame
    ctf.text = "PEAT-Nucleate BioHack 2026  •  HistoTME + TCGA Lung"
    ctf.paragraphs[0].font.size = Pt(14)
    ctf.paragraphs[0].font.color.rgb = TEAL
    _add_notes(
        s,
        "Open with the core tension: trials stratify by mutation, but patients with the same "
        "EGFR, ALK, or KRAS alteration respond very differently. The missing variable is the "
        "tumor microenvironment and how immune and stromal cells are organized in space.",
    )

    # --- Slide 2: The problem ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, WHITE)
    _add_title_bar(s, "The clinical problem", "Mutation panels are necessary — but not sufficient")
    _add_bullets(
        s,
        [
            "NSCLC trials increasingly enroll by driver: EGFR, ALK, KRAS G12C, MET, RET…",
            "Initial targeted responses can be strong, yet median benefit is limited by resistance",
            "Immunotherapy combinations often fail in oncogene-driven disease despite 'targetable' biology",
            "Trial readouts report averages for 'EGFR+' or 'KRAS G12C+' — masking within-group heterogeneity",
            "",
            "Key question: What differs between responders and non-responders with the same mutation?",
        ],
        font_size=19,
    )
    _add_footer(s, footer)
    _add_notes(
        s,
        "Emphasize that this is not about finding new mutations. The mutation is known and druggable. "
        "The puzzle is why two EGFR-mutant patients on the same TKI or trial arm have different durability.",
    )

    # --- Slide 3: Trial examples ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, WHITE)
    _add_title_bar(s, "Trials with the right mutation — disappointing results")
    _add_table_slide(
        s,
        ["Trial", "Population", "Intervention", "Outcome"],
        [
            ["KEYNOTE-789", "EGFR+ post-TKI", "Pembro + chemo vs chemo", "No significant PFS/OS benefit"],
            ["CheckMate 722", "EGFR+ post-TKI", "Nivo + chemo vs chemo", "No significant benefit"],
            ["TATTON", "EGFR-mutant", "Osimertinib + durvalumab", "Arm halted — ILD toxicity"],
            ["IMpower151", "EGFR/ALK enriched", "Atezo + bev + chemo", "Negative primary PFS endpoint"],
            ["CodeBreaK 200", "KRAS G12C", "Sotorasib vs docetaxel", "PFS improved; OS not improved"],
        ],
        [1.35, 1.55, 2.35, 2.85],
    )
    _add_footer(s, footer)
    _add_notes(
        s,
        "These are anchor examples for judges and clinicians. The through-line: genotype-qualified "
        "patients still fail or see modest benefit. Resistance mechanisms include bypass signaling, "
        "but also immunosuppressive TME remodeling.",
    )

    # --- Slide 4: Hypothesis ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, WHITE)
    _add_title_bar(s, "Hypothesis", "TME architecture explains residual heterogeneity")
    _add_bullets(
        s,
        [
            "Patients with identical driver mutations occupy distinct tumor microenvironment states",
            "  Immune desert vs immune inflamed",
            "  Macrophage/Treg-rich niches vs effector T-cell infiltration",
            "  CAF/stromal barriers vs angiogenic remodeling",
            "",
            "Spatial organization matters — not just bulk immune scores",
            "Genomics tells us what to target; TME tells us whether therapy will work",
            "",
            "Spatial transcriptomics is gold standard — but costly and rarely run at trial scale",
            "We need a scalable first pass on standard diagnostic H&E",
        ],
        font_size=18,
        level_indent=True,
    )
    _add_footer(s, footer)
    _add_notes(
        s,
        "Clarify nuance: EGFR/ALK tumors often have LOW tumor mutational burden and cold immune "
        "phenotypes — the issue is not always 'too many mutations.' For KRAS G12C, resistance "
        "includes TME shifts toward immunologically cold states after treatment.",
    )

    # --- Slide 5: What is HistoTME ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, WHITE)
    _add_title_bar(s, "Our approach: HistoTME", "Infer TME programs from H&E whole-slide images")
    img_path = ROOT / "external" / "HistoTME" / "figures" / "HistoTME_outline.png"
    if img_path.exists():
        s.shapes.add_picture(str(img_path), Inches(0.45), Inches(1.35), width=Inches(5.5))
    _add_bullets(
        s,
        [
            "Weakly supervised deep learning on pathology foundation-model embeddings",
            "Predicts 29 curated gene-expression signatures",
            "No extra staining; works on archived diagnostic slides",
            "Published in NSCLC; pan-cancer HistoTMEv2 extension",
            "",
            "Bulk mode → patient/slide-level TME profile",
            "Spatial mode → tile-level maps across the WSI",
        ],
        left=6.1,
        top=1.45,
        width=3.5,
        font_size=15,
    )
    _add_footer(s, footer)
    _add_notes(
        s,
        "HistoTME bridges histology and transcriptomics. It was originally developed for NSCLC "
        "immunotherapy response prediction. HistoTMEv2 extends across 25 cancer types including lung.",
    )

    # --- Slide 6: 29 signatures ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, WHITE)
    _add_title_bar(s, "29 interpretable TME signatures", "Biology clinicians recognize — not a black-box score")
    _add_bullets(
        s,
        [
            "Immune-active: Effector_cells, T_cells, NK_cells, M1_signatures, Antitumor_cytokines, MHCII",
            "Immunosuppressive: Treg, MDSC, Checkpoint_inhibition, Th2_signature, Protumor_cytokines",
            "Myeloid: Macrophages, Macrophage_DC_traffic, Neutrophil_signature, Granulocyte_traffic",
            "Stroma / escape: CAF, Angiogenesis, Matrix_remodeling, EMT_signature",
            "Tumor state: Proliferation_rate",
            "",
            "Downstream clustering → Immune Desert vs Immune Inflamed archetypes",
            "Validated against IHC (CD4/CD8, B cells, macrophages) in published analyses",
        ],
        font_size=17,
    )
    _add_footer(s, footer)
    _add_notes(
        s,
        "These signatures come from curated bulk RNA programs. The downstream module shows "
        "Immune Inflamed phenotype predicts better ICI outcomes than PD-L1 alone in several subgroups.",
    )

    # --- Slide 7: Inference modes figure ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, WHITE)
    _add_title_bar(s, "Bulk vs spatial inference", "One slide → whole-patient profile or spatial TME map")
    img_path = ROOT / "external" / "HistoTME" / "figures" / "inference_modes_figure.png"
    if img_path.exists():
        s.shapes.add_picture(str(img_path), Inches(0.55), Inches(1.3), height=Inches(5.6))
    _add_footer(s, footer)
    _add_notes(
        s,
        "Spatial mode is the differentiator for your pitch. predict_spatial.py outputs enrichment "
        "scores per tile coordinate — a computational analogue of spatial transcriptomics screening.",
    )

    # --- Slide 8: Repo pipeline ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, WHITE)
    _add_title_bar(s, "PEAT-Nucleate pipeline (this repo)", "TCGA lung cohort ready to run")
    _add_bullets(
        s,
        [
            "1,053 diagnostic H&E slides — TCGA-LUAD (541) + TCGA-LUSC (512)",
            "  ~824 GB total; open access; manifests committed in data/tcga_lung/",
            "",
            "Pipeline:",
            "  download.py → whole-slide images (.svs)",
            "  Foundation model embeddings → HDF5 (coords + features)",
            "  predict_bulk.py / predict_spatial.py → TME signature outputs",
            "  Link to clinical + mutation metadata for stratified analysis",
            "",
            "Outputs land in outputs/histotme/ — ready for dashboard (Taylor UI branch)",
        ],
        font_size=17,
        level_indent=True,
    )
    _add_footer(s, footer)
    _add_notes(
        s,
        "Walk through the repo as an end-to-end demo path. Pilot with --limit 3 slides for the hackathon. "
        "Clinical metadata branch adds patient-level joins for EGFR/KRAS/ALK stratification.",
    )

    # --- Slide 9: Positioning vs spatial transcriptomics ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, WHITE)
    _add_title_bar(s, "HistoTME + spatial transcriptomics", "Complementary — not competing")
    _add_table_slide(
        s,
        ["", "Spatial transcriptomics", "HistoTME on H&E"],
        [
            ["Input", "Fresh/frozen; specialized workflow", "Standard diagnostic H&E"],
            ["Scale", "Limited spots / FOVs", "Whole slide; 1,000+ TCGA lungs"],
            ["Cost & time", "High", "Compute-only after slide download"],
            ["Output", "Measured genes at locations", "Predicted signature activity per tile"],
            ["Best role", "Mechanistic ground truth", "Trial-scale screening & stratification"],
        ],
        [1.6, 3.2, 3.3],
    )
    quote = s.shapes.add_textbox(Inches(0.5), Inches(6.55), Inches(9), Inches(0.55))
    qtf = quote.text_frame
    qtf.text = "Workflow: HistoTME screens the cohort → spatial transcriptomics validates on a focused subset"
    qtf.paragraphs[0].font.size = Pt(13)
    qtf.paragraphs[0].font.italic = True
    qtf.paragraphs[0].font.color.rgb = TEAL
    _add_footer(s, footer)
    _add_notes(
        s,
        "This slide prevents the 'but it's not real spatial omics' objection. You're proposing a "
        "two-layer strategy: scalable H&E inference first, spatial transcriptomics validation second.",
    )

    # --- Slide 10: Proposed analysis ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, WHITE)
    _add_title_bar(s, "Proposed analysis", "Test the hypothesis on TCGA lung")
    _add_bullets(
        s,
        [
            "Stratify patients by driver mutation (EGFR, ALK, KRAS, wild-type)",
            "Run HistoTME bulk + spatial on diagnostic H&E slides",
            "Cluster into TME archetypes (immune desert, inflamed, stroma-high, myeloid-rich)",
            "Compare TME states within each mutation stratum",
            "",
            "Retrospective trial framing:",
            "  Do non-responders in mutation-positive trials show distinct spatial TME maps?",
            "  Can TME archetype explain heterogeneity that mutation status alone cannot?",
            "",
            "Validate signature maps against public spatial transcriptomics lung datasets",
        ],
        font_size=17,
        level_indent=True,
    )
    _add_footer(s, footer)
    _add_notes(
        s,
        "TCGA lacks trial treatment arms, so you frame this as retrospective hypothesis generation "
        "that motivates prospective biopsy studies and trial enrichment designs.",
    )

    # --- Slide 11: BioHack deliverables ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, WHITE)
    _add_title_bar(s, "BioHack deliverables", "What we will show live")
    _add_bullets(
        s,
        [
            "Pilot TCGA lung download (3–10 slides) + tissue thumbnails",
            "HistoTME bulk predictions for LUAD/LUSC pilot cohort",
            "Spatial heatmaps for 2–3 exemplar slides (Treg, Macrophages, Effector_cells)",
            "Mutation-stratified TME cluster comparison (EGFR vs KRAS vs WT)",
            "Dashboard mock-up: slide browser + TME overlay (Taylor UI branch)",
            "",
            "Stretch goal: correlate HistoTME signatures with spatial transcriptomics public data",
        ],
        font_size=18,
    )
    _add_footer(s, footer)
    _add_notes(
        s,
        "Be concrete about what is demo-ready vs stretch. Judges reward a working pilot over promises.",
    )

    # --- Slide 12: Impact ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, NAVY)
    impact = s.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(8.5), Inches(4.5))
    itf = impact.text_frame
    lines = [
        "Impact",
        "",
        "• Enrich future targeted-therapy and combination trials by TME — not mutation alone",
        "• Explain post-hoc why mutation-positive trials underperform in subsets",
        "• Use existing H&E biobanks — no new biopsy protocol required for screening",
        "• Bridge to spatial transcriptomics for mechanistic validation",
        "",
        "Mutation is the lock. TME is whether the key turns.",
    ]
    for i, line in enumerate(lines):
        p = itf.paragraphs[0] if i == 0 else itf.add_paragraph()
        p.text = line
        if i == 0:
            p.font.size = Pt(36)
            p.font.bold = True
        elif line.startswith("Mutation is"):
            p.font.size = Pt(22)
            p.font.italic = True
            p.font.color.rgb = TEAL
        elif line == "":
            p.font.size = Pt(6)
        else:
            p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)
    _add_notes(
        s,
        "Close with the clinical translation path: retrospective TCGA → trial biobank re-analysis → "
        "prospective enrichment biomarker on pre-treatment H&E.",
    )

    # --- Slide 13: References ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, WHITE)
    _add_title_bar(s, "Key references")
    _add_bullets(
        s,
        [
            "Patkar et al. NPJ Precision Oncology 2024 — HistoTME NSCLC / ICI",
            "HistoTMEv2 preprint (2025) — pan-cancer extension",
            "Yang et al. J Clin Oncol 2024 — KEYNOTE-789 (EGFR+ post-TKI)",
            "de Langen et al. Lancet 2023 — CodeBreaK 200 (KRAS G12C)",
            "Tsai et al. JCI — KRAS G12Ci resistance & TME remodeling",
            "Nature Reviews Clinical Oncology 2025 — NSCLC TME crosstalk",
            "Repo: github.com/aonkondey01/PEAT-Nucleate-BIoHack-2026",
        ],
        font_size=15,
    )
    _add_footer(s, footer)
    _add_notes(s, "Keep this slide for Q&A. Offer to share the deck and pilot notebooks after the talk.")

    # --- Slide 14: Thank you ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, TEAL)
    tbox = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = tbox.text_frame
    tf.text = "Thank you"
    tf.paragraphs[0].font.size = Pt(48)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    q = s.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1.2))
    qtf = q.text_frame
    qtf.text = "Questions?\nPEAT-Nucleate BioHack 2026"
    for i, p in enumerate(qtf.paragraphs):
        p.font.size = Pt(20 if i == 0 else 16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    _add_notes(s, "Invite questions on trial selection, spatial validation plan, and live demo timeline.")

    prs.save(OUT_PATH)
    return OUT_PATH


if __name__ == "__main__":
    path = build_deck()
    print(f"Wrote {path} ({path.stat().st_size // 1024} KB)")
